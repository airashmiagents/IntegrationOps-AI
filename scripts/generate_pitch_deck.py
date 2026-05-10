#!/usr/bin/env python3
"""
Generate IntegrationOps AI hackathon deck (dark / enterprise theme).

Run from repo root:
  pip install -r scripts/requirements-pptx.txt && python scripts/generate_pitch_deck.py

Output (gitignored): presentation/IntegrationOps-AI-Deck.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# --- Theme (slate + indigo / cyan enterprise) ---
BG = RGBColor(15, 23, 42)  # #0f172a
BG_CARD = RGBColor(30, 41, 59)  # #1e293b
TEXT = RGBColor(248, 250, 252)  # #f8fafc
TEXT_MUTED = RGBColor(148, 163, 184)  # #94a3b8
ACCENT_A = RGBColor(99, 102, 241)  # indigo-500
ACCENT_B = RGBColor(56, 189, 248)  # sky-400
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
_REPO_ROOT = Path(__file__).resolve().parent.parent
_OUT_DIR = _REPO_ROOT / "presentation"
_OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT = _OUT_DIR / "IntegrationOps-AI-Deck.pptx"


def _full_bleed_bg(slide, color: RGBColor = BG) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _gradient_header(slide) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.22))
    fill = bar.fill
    fill.gradient()
    fill.gradient_angle = 0.0
    fill.gradient_stops[0].color.rgb = ACCENT_A
    fill.gradient_stops[1].color.rgb = ACCENT_B
    bar.line.fill.background()


def _slide_label(slide, label: str) -> None:
    box = slide.shapes.add_textbox(Inches(11.85), Inches(7.05), Inches(1.4), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.RIGHT


def _title(slide, text: str, top: float = 0.55, size: int = 40) -> None:
    box = slide.shapes.add_textbox(Inches(0.65), Inches(top), Inches(12), Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = TEXT
    p.font.name = "Segoe UI"


def _subtitle(slide, text: str, top: float = 1.35, size: int = 16) -> None:
    box = slide.shapes.add_textbox(Inches(0.65), Inches(top), Inches(12), Inches(0.55))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = TEXT_MUTED
    p.font.name = "Segoe UI Light"


def _icon_badge(slide, emoji: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.65), Inches(0.28), Inches(0.55), Inches(0.45))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = emoji
    p.font.size = Pt(28)


def _bullets(slide, lines: list[str], top: float = 2.05, left: float = 0.85) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(11.6), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20 if len(lines) <= 4 else 17)
        p.font.color.rgb = TEXT
        p.font.name = "Segoe UI"
        p.space_after = Pt(14)
        p.level = 0


def _flow_architecture(slide) -> None:
    """Two-row pipeline: ingest → intelligence → outcomes."""
    y1 = Inches(2.05)
    y2 = Inches(2.92)
    w1, h1 = Inches(1.38), Inches(0.58)
    w2, h2 = Inches(1.85), Inches(0.52)
    gap = Inches(0.1)
    row1 = [
        "SAP CPI",
        "API layer",
        "AI agent",
        "Context builder",
        "LLM",
        "RCA",
    ]
    row2 = ["Jira", "Email", "Dashboard"]
    x0 = Inches(0.38)

    def row(items: list[str], y: float, w: float, h: float) -> None:
        for i, label in enumerate(items):
            x = x0 + i * (w + gap)
            sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            sh.fill.solid()
            sh.fill.fore_color.rgb = BG_CARD
            sh.line.color.rgb = ACCENT_B
            sh.line.width = Pt(1)
            tf = sh.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(10 if w < Inches(1.5) else 11)
            p.font.bold = True
            p.font.color.rgb = TEXT
            p.alignment = PP_ALIGN.CENTER

    row(row1, y1, w1, h1)
    row2_width = len(row2) * w2 + (len(row2) - 1) * gap
    x2 = (SLIDE_W - row2_width) / 2
    for j, label in enumerate(row2):
        x = x2 + j * (w2 + gap)
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y2, w2, h2)
        sh.fill.solid()
        sh.fill.fore_color.rgb = BG_CARD
        sh.line.color.rgb = ACCENT_A
        sh.line.width = Pt(1)
        tf = sh.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT
        p.alignment = PP_ALIGN.CENTER

    caption = (
        "OData MPL + design-time metadata → FastAPI → run_investigation → build_context → "
        "OpenRouter (DeepSeek / fallback) → structured RCA → SQLite + mock inbox + Jira-ready fields"
    )
    cap_box = slide.shapes.add_textbox(Inches(0.45), Inches(3.05), Inches(12.2), Inches(1.2))
    tf = cap_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = caption
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_MUTED
    p.font.name = "Segoe UI"


def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # --- 1 Title ---
    s = prs.slides.add_slide(blank)
    _full_bleed_bg(s)
    _gradient_header(s)
    _icon_badge(s, "◆")
    _title(s, "IntegrationOps AI", top=1.55, size=48)
    _subtitle(s, "Autonomous SAP CPI Incident Intelligence Platform", top=2.45, size=20)
    tag = s.shapes.add_textbox(Inches(0.65), Inches(3.35), Inches(11), Inches(0.5))
    tp = tag.text_frame.paragraphs[0]
    tp.text = "Hackathon build · Enterprise reliability · AI-first operations"
    tp.font.size = Pt(14)
    tp.font.color.rgb = ACCENT_B
    tp.font.name = "Segoe UI"
    _slide_label(s, "01 / Title")

    # --- 2 Problem ---
    s = prs.slides.add_slide(blank)
    _full_bleed_bg(s)
    _gradient_header(s)
    _icon_badge(s, "⚡")
    _title(s, "Problem", size=36)
    _subtitle(s, "Integration failures drain time and trust", top=1.22)
    _bullets(
        s,
        [
            "SAP CPI Message Processing Logs are dense — RCA is expert-led and slow.",
            "Manual correlation across MPL, design-time artifacts, and tickets is error-prone.",
            "Mean time to resolve stays high when context is fragmented across tools.",
        ],
        top=2.0,
    )
    _slide_label(s, "02 / Problem")

    # --- 3 Solution ---
    s = prs.slides.add_slide(blank)
    _full_bleed_bg(s)
    _gradient_header(s)
    _icon_badge(s, "◇")
    _title(s, "Solution", size=36)
    _subtitle(s, "One autonomous agent, one structured outcome", top=1.22)
    _bullets(
        s,
        [
            "Autonomous investigation agent — same pipeline as operator “deep dive”.",
            "Pulls CPI runtime logs + iFlow metadata (OData) into a single LLM briefing.",
            "Returns RCA, severity, confidence, and recommendations as enterprise JSON.",
        ],
        top=2.0,
    )
    _slide_label(s, "03 / Solution")

    # --- 4 Architecture ---
    s = prs.slides.add_slide(blank)
    _full_bleed_bg(s)
    _gradient_header(s)
    _icon_badge(s, "⬡")
    _title(s, "Architecture", size=36)
    _subtitle(s, "End-to-end signal → decision → action surface", top=1.22)
    _flow_architecture(s)
    _slide_label(s, "04 / Architecture")

    # --- 5 Features ---
    s = prs.slides.add_slide(blank)
    _full_bleed_bg(s)
    _gradient_header(s)
    _icon_badge(s, "✦")
    _title(s, "Key Features", size=36)
    _subtitle(s, "Shipped in the MVP + clear hooks for scale", top=1.22)
    _bullets(
        s,
        [
            "◎ Real-time / scheduled CPI FAILED-MPL monitor (APScheduler) with dedupe.",
            "◎ AI-driven RCA + confidence scoring (OpenRouter + heuristic fallback).",
            "◎ Persisted incidents (SQLite) + observability lifecycle (LLM audit join).",
            "◎ Jira-ready fields + mock alert inbox (email preview without SMTP).",
            "◎ Operator dashboard + “run monitor now” for instant demos.",
        ],
        top=1.95,
        left=0.75,
    )
    _slide_label(s, "05 / Features")

    # --- 6 Demo flow ---
    s = prs.slides.add_slide(blank)
    _full_bleed_bg(s)
    _gradient_header(s)
    _icon_badge(s, "▶")
    _title(s, "Live Demo Flow", size=36)
    _subtitle(s, "What judges see in under two minutes", top=1.22)
    _bullets(
        s,
        [
            "① Incident detected — FAILED MPL in lookback window.",
            "② Logs + metadata fetched — CPI OData into the agent.",
            "③ AI analysis — LLM JSON with evidence-aligned confidence.",
            "④ RCA generated — dashboard + mock inbox + SQLite audit trail.",
            "⑤ Jira + email — schema-ready / mock UI today; wire to ITSM tomorrow.",
        ],
        top=1.95,
    )
    _slide_label(s, "06 / Demo")

    # --- 7 Tech stack ---
    s = prs.slides.add_slide(blank)
    _full_bleed_bg(s)
    _gradient_header(s)
    _icon_badge(s, "⌁")
    _title(s, "Tech Stack", size=36)
    _subtitle(s, "Modern, boring-in-a-good-way production defaults", top=1.22)
    _bullets(
        s,
        [
            "FastAPI · Uvicorn · Pydantic — typed APIs, OpenAPI docs, /api prefix for proxies.",
            "React 18 + Vite + React Router — dashboard, lifecycle, mock inbox.",
            "OpenRouter (DeepSeek + Llama fallback) — httpx, JSON-mode prompts.",
            "SAP CPI OData (MPL + Integration Content) — mock + live paths.",
            "Render (API) + Vercel (UI) — env-driven VITE_API_URL, CORS-ready.",
        ],
        top=1.9,
        left=0.75,
    )
    _slide_label(s, "07 / Stack")

    # --- 8 Impact ---
    s = prs.slides.add_slide(blank)
    _full_bleed_bg(s)
    _gradient_header(s)
    _icon_badge(s, "↑")
    _title(s, "Impact", size=36)
    _subtitle(s, "Why operators and leadership care", top=1.22)
    _bullets(
        s,
        [
            "↓ MTTR — first-pass triage with cited evidence, not tribal knowledge.",
            "↑ Consistency — same agent steps every incident; auditable LLM exchanges.",
            "↑ Reliability posture — confidence scores + severity for prioritization.",
        ],
        top=2.1,
    )
    _slide_label(s, "08 / Impact")

    # --- 9 Future ---
    s = prs.slides.add_slide(blank)
    _full_bleed_bg(s)
    _gradient_header(s)
    _icon_badge(s, "→")
    _title(s, "Future Scope", size=36)
    _subtitle(s, "From insight to autonomous remediation", top=1.22)
    _bullets(
        s,
        [
            "Auto-healing playbooks — safe rollbacks, cache clears, credential rotation hooks.",
            "Event-driven remediation — CPI + ITSM webhooks feeding a closed loop.",
            "Multi-system observability — extend the agent beyond CPI (IDoc, BTP, A2A).",
        ],
        top=2.05,
    )
    _slide_label(s, "09 / Future")

    # --- 10 Thank you / CTA ---
    s = prs.slides.add_slide(blank)
    _full_bleed_bg(s)
    _gradient_header(s)
    _icon_badge(s, "◆")
    _title(s, "IntegrationOps AI", top=2.35, size=44)
    _subtitle(s, "Autonomous SAP CPI Incident Intelligence", top=3.25, size=18)
    cta = s.shapes.add_textbox(Inches(0.65), Inches(4.15), Inches(12), Inches(0.6))
    cp = cta.text_frame.paragraphs[0]
    cp.text = "Live API · Dashboard · Mock inbox · Observability drill-down"
    cp.font.size = Pt(16)
    cp.font.color.rgb = ACCENT_B
    cp.font.name = "Segoe UI"
    _slide_label(s, "10 / Close")

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
